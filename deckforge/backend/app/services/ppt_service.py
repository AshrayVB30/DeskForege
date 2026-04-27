from pptx import Presentation
import io
import httpx
from pptx.util import Inches

def create_ppt_from_json(slides_data: list) -> io.BytesIO:
    """
    Takes JSON structure and creates a PPTX binary stream.
    Expected format: [{"title": "str", "points": ["str"], "image_url": "str"}]
    """
    prs = Presentation()
    
    # 0 is title slide layout, 1 is title and content layout
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Optional: Title slide generated from first item
    if len(slides_data) > 0:
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = slides_data[0].get("title", "Presentation")
        points = slides_data[0].get("points", [])
        if points:
            subtitle.text = "\n".join(points)
            
    # Rest of the slides
    for slide_data in slides_data[1:]:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_data.get("title", "Untitled Slide")
        
        # Adjust body shape for image if present
        image_url = slide_data.get("image_url")
        if image_url:
            # Resize text box to half width
            body_shape.width = Inches(4.5)
            
            try:
                # Download image
                with httpx.Client() as client:
                    img_resp = client.get(image_url)
                    if img_resp.status_code == 200:
                        img_data = io.BytesIO(img_resp.content)
                        # Add image to the right side
                        slide.shapes.add_picture(img_data, Inches(5), Inches(1.5), width=Inches(4))
            except Exception as e:
                print(f"Error adding image to PPT: {e}")

        tf = body_shape.text_frame
        points = slide_data.get("points", [])
        
        if points:
            tf.text = points[0]
            for point in points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream
